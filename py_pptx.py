import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import openpyxl
from pptx.enum.chart import XL_CHART_TYPE
import openpyxl
import os
# import win32com.client as win32
from pptx import Presentation
from pptx.util import Inches

Image1 = "/home/hunter/Documents/Workspace/Python_code/py_pptx/image1.jpg"
Image2 = "/home/hunter/Documents/Workspace/Python_code/py_pptx/image2.jpg"

def generate_ppt(image1_path, image2_path, title,pres_name):

    prs = Presentation()

    # set width and height to 16 and 9 inches.
    slide_width  = prs.slide_width = Inches(16)
    slide_height =prs.slide_height = Inches(9)

    # create first slide with image1
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    pic1 = slide1.shapes.add_picture(image1_path, 0, 0, slide_width, slide_height)
    
    # create second slide with image2
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    pic2 = slide2.shapes.add_picture(image2_path, 0, 0, slide_width, slide_height)

    # add title to both slides
    # for slide in [slide1, slide2]:
    title_shape = slide2.shapes.add_textbox(Inches(2), Inches(5.5), Inches(4), Inches(0.5))
    title_text_frame = title_shape.text_frame
    title_text_frame.text = title
    title_text_frame.paragraphs[0].font.bold = True
    title_text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_text_frame.paragraphs[0].font.size = Pt(48)
    title_shape.top = prs.slide_height - title_shape.height - Inches(2.2)

    # Define background image path
    bg_image_path = 'Background.jpg'

    # Loop through and add 12 slides
    for i in range(1, 13):
        # Add slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Add background image
        slide.background.fill.solid()
        slide.background.fill.background()
        pic = slide.shapes.add_picture(bg_image_path, 0, 0, prs.slide_width, prs.slide_height)

        # Add header text box with title
        header_shape = slide.shapes.add_textbox(Inches(0), Inches(0.5), Inches(9), Inches(1))
        header_text_frame = header_shape.text_frame
        header_text_frame.text = f'Slide {i} Title'
        header_text_frame.paragraphs[0].font.size = Inches(0.5)
        header_text_frame.paragraphs[0].font.bold = True
        header_text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        # Add body text box with small title
        body_shape = slide.shapes.add_textbox(Inches(0), Inches(1.5), Inches(9), Inches(1.5))
        body_text_frame = body_shape.text_frame
        body_text_frame.text = f'pptx_file (str): The file path of the \n PowerPoint file to which to append the slides.'
        body_text_frame.paragraphs[0].font.size = Inches(0.3)
        body_text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        body_text_frame.paragraphs[0].font.bold = True
        # Adjust left position of text boxes to create indentation
        header_shape.left = Inches(0)
        body_shape.left = Inches(0)
    
    slide_width  = prs.slide_width = Inches(16)
    slide_height =prs.slide_height = Inches(9)
    # create second slide with image2
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    pic2 = slide2.shapes.add_picture(Image1, 0, 0, slide_width, slide_height)
    prs.save(f"{pres_name}.pptx")

    #!-----------------------------------------------------------------------------


    # Load the Excel workbook.
    workbook = openpyxl.load_workbook(filename='step_chart.xlsx')

    # Load the PowerPoint presentation.
    ppt = Presentation('presentation.pptx')

    # Loop through each worksheet in the Excel workbook
    for sheet in workbook:

        # Loop through each chart in the worksheet
        for chart in sheet._charts:

            # Create a new slide in the PowerPoint presentation
            slide = ppt.slides.add_slide(ppt.slide_layouts[6])

            # Add the chart to the slide
            chart_blob = chart._write._chart_part()
            chart_object = slide.shapes.add_chart(chart.chart_type, Inches(1), Inches(2), Inches(8), Inches(4), chart_blob)


        # Save the modified PowerPoint presentation.
        ppt.save('presentation.pptx')


    


    #!...............................................................
generate_ppt(Image1, Image2, 'My Presentation Title',"presentation")
